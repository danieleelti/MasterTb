import streamlit as st
import streamlit.components.v1 as components
import gspread
import pandas as pd
import json
import re
import ast
import io
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import pypdf
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import difflib

# --- CONFIGURAZIONE PAJINA (WIDE MODE) ---
st.set_page_config(layout="wide", page_title="MasterTb Manager", page_icon="🦁")

# --- CONFIGURAZIONE MODELLI FISSI ---
SEARCH_MODEL = "models/gemini-2.5-flash-lite"
DOC_MODEL = "models/gemini-3-pro-preview"

# --- CSS E JS PER DRAG & DROP VISUALE ---
components.html("""
<script>
    const observer = new MutationObserver(() => {
        const dropzone = window.parent.document.querySelector('[data-testid="stFileUploaderDropzone"]');
        if (dropzone && !dropzone.classList.contains('drag-listener')) {
            dropzone.classList.add('drag-listener');
            dropzone.addEventListener('dragover', (e) => {
                e.preventDefault();
                dropzone.style.border = '3px dashed #FF4B4B';
                dropzone.style.backgroundColor = '#FFECEC';
                dropzone.style.transform = 'scale(1.02)';
            });
            const reset = () => {
                dropzone.style.border = '1px dashed #aaa';
                dropzone.style.backgroundColor = '#f9f9f9';
                dropzone.style.transform = 'scale(1.0)';
            };
            dropzone.addEventListener('dragleave', reset);
            dropzone.addEventListener('drop', reset);
        }
    });
    observer.observe(window.parent.document.body, { childList: true, subtree: true });
</script>
""", height=0, width=0)

st.markdown("""
    <style>
    [data-testid='stFileUploaderDropzone'] {
        border: 1px dashed #aaa;
        border-radius: 10px;
        background-color: #f9f9f9;
        transition: all 0.2s ease-in-out;
        padding: 10px;
    }
    .diff-box {
        padding: 10px;
        border-radius: 5px;
        margin-bottom: 5px;
        border: 1px solid #ddd;
    }
    .old-val { background-color: #ffe6e6; color: #b30000; text-decoration: line-through; padding: 2px 5px; border-radius: 3px;}
    .new-val { background-color: #e6fffa; color: #006644; font-weight: bold; padding: 2px 5px; border-radius: 3px;}
    
    .stTextArea textarea[aria-label="Modifica Nuova Descrizione"] {
        background-color: #e6fffa;
        border: 2px solid #006644;
        color: #004d33;
    }
    
    /* Stile Card Risultati in Sidebar */
    .result-card {
        background-color: white;
        padding: 10px;
        border-radius: 5px;
        border: 1px solid #eee;
        margin-bottom: 8px;
        box-shadow: 0 1px 2px rgba(0,0,0,0.05);
    }
    .result-title { font-weight: bold; font-size: 0.9em; margin-bottom: 4px; }
    .result-preview { font-size: 0.8em; color: #666; margin-bottom: 8px; line-height: 1.2; }
    </style>
""", unsafe_allow_html=True)

# --- INIZIALIZZAZIONE STATO ---
if 'logged_in' not in st.session_state: st.session_state['logged_in'] = False
if 'token_usage' not in st.session_state: 
    st.session_state['token_usage'] = {'input': 0, 'output': 0, 'total': 0}
if 'search_results' not in st.session_state:
    st.session_state['search_results'] = None
if 'pending_duplicate' not in st.session_state:
    st.session_state['pending_duplicate'] = None
if 'last_processed_file' not in st.session_state:
    st.session_state['last_processed_file'] = None
if 'force_selection' not in st.session_state:
    st.session_state['force_selection'] = None

# Stato per modifiche in attesa di conferma
if 'pending_changes' not in st.session_state:
    st.session_state['pending_changes'] = None

# Inizializzazione sicura draft_data (Per nuovi upload)
if 'draft_data' not in st.session_state:
    st.session_state['draft_data'] = {}
elif st.session_state['draft_data'] is None:
    st.session_state['draft_data'] = {}

# Variabili Debug
if 'debug_raw_text' not in st.session_state: st.session_state['debug_raw_text'] = ""
if 'debug_ai_response' not in st.session_state: st.session_state['debug_ai_response'] = ""

# --- 1. LOGIN ---
if not st.session_state['logged_in']:
    st.title("🦁 MasterTb Accesso")
    with st.form("login"):
        pwd = st.text_input("Password", type="password")
        if st.form_submit_button("Entra"):
            if "login_password" in st.secrets and pwd == st.secrets["login_password"]:
                st.session_state['logged_in'] = True
                st.rerun()
            else:
                st.error("Password errata")
    st.stop()

# --- 2. CONNESSIONE ---
ws = None
@st.cache_resource
def connect_to_sheet():
    try:
        creds = st.secrets["gcp_service_account"]
        gc = gspread.service_account_from_dict(creds)
        return gc.open("MasterTbGoogleAi").get_worksheet(0)
    except Exception as e:
        st.error(f"Errore Sheet: {e}")
        st.stop()

ws = connect_to_sheet()

@st.cache_data(ttl=60)
def load_data():
    df = pd.DataFrame(ws.get_all_records())
    if not df.empty:
        df.columns = [c.strip() for c in df.columns]
        df.set_index(df.columns[0], inplace=True)
    return df

df = load_data()
if df.empty: st.stop()

product_ids = [str(i) for i in df.index.tolist()]
cols = df.columns.tolist()
id_col = df.index.name

# --- HELPER UTILITY ---
def get_shape_text_recursive(shape):
    text = ""
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            text += shape.text_frame.text + "\n"
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        text += cell.text_frame.text + " "
            text += "\n"
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                text += get_shape_text_recursive(child)
    except: pass
    return text

def read_file_content(uploaded_file):
    text = ""
    try:
        if uploaded_file.name.endswith('.pdf'):
            pdf_reader = pypdf.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        elif uploaded_file.name.endswith('.pptx') or uploaded_file.name.endswith('.ppt'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    text += get_shape_text_recursive(shape)
    except Exception as e:
        st.error(f"Errore lettura file: {e}")
    return text

def create_slug(text):
    if not text: return ""
    text = text.lower().strip()
    text = re.sub(r'[^a-z0-9\s-]', '', text)
    text = re.sub(r'\s+', '-', text)
    return text

# --- 3. FUNZIONI AI ---
def analyze_document_with_gemini(text_content, columns):
    if "GOOGLE_API_KEY" not in st.secrets: return {}
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])

    # Identifica colonne chiave
    desc_col_name = "Descrizione Breve"
    log_col_name = "Logistica"
    
    for c in columns:
        if "descrizione" in c.lower(): desc_col_name = c
        if "logistica" in c.lower(): log_col_name = c

    # Prompt SPECIFICO per Esperto Team Building
    sys_prompt = f"""
    Sei un SENIOR TEAM BUILDING EXPERT. Analizza il documento fornito per estrarre informazioni strategiche.
    
    OBIETTIVO: Compilare un JSON con queste chiavi esatte:
    {json.dumps(columns)}
    
    1. Campo Chiave: '{columns[0]}' (NOME FORMAT).
    2. Campo '{desc_col_name}': SCRIVI 5-6 RIGHE COMPLETE coinvolgenti e descrittive.
    3. Campo '{log_col_name}': Estrai dettagli tecnici, spazi (indoor/outdoor), necessità (tavoli, corrente, acqua). Sii preciso.
    
    REGOLE DI RAGIONAMENTO (THINKING PROCESS):
    - 'Target Ideale': NON copiare solo il testo. Ragiona: a chi si rivolge? Sales? Management? Tutti? Scrivi una sintesi mirata.
    - 'Formazione': Analizza se l'attività sviluppa Soft Skills (Leadership, Comunicazione, Problem Solving). Se sì, descrivile brevemente. Se è solo ludico, scrivi "Ludico/Incentive".
    - 'Sociale': Analizza se l'attività spinge forte sull'interazione e condivisione. Rispondi SOLO "SI" o "NO".
    - 'Ranking': Valuta la complessità logistica e l'impatto emotivo da 1 a 5 basandoti sulla tua esperienza. Rispondi SOLO col numero.
    - 'Durata' (Min/Max/Media): Analizza i tempi. Se trovi un range, calcola tu la MEDIA.
    - 'Max Pax': Se non trovi un limite specifico, scrivi "illimitato".
    - 'Metodo di Calcolo': Se non specificato diversamente, ipotizza "Standard".
    
    REGOLE FORMALI:
    - Se l'informazione MANCA DEL TUTTO, scrivi "[[RIEMPIMENTO MANUALE]]".
    - Rispondi SOLO con il JSON.
    """

    model = genai.GenerativeModel(
        model_name=DOC_MODEL,
        generation_config={"temperature": 0.2, "response_mime_type": "application/json"},
        system_instruction=sys_prompt
    )

    try:
        response = model.generate_content(f"TESTO DOCUMENTO:\n{text_content}")
        clean_text = response.text.strip()
        if clean_text.startswith("```json"): clean_text = clean_text[7:]
        if clean_text.endswith("```"): clean_text = clean_text[:-3]
        
        st.session_state['debug_ai_response'] = clean_text 
        return json.loads(clean_text.strip())
    except Exception as e:
        st.error(f"Errore AI ({DOC_MODEL}): {e}")
        return {}

def search_ai(query, dataframe):
    if "GOOGLE_API_KEY" not in st.secrets: return []
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
    
    context_str = dataframe.to_markdown(index=True)
    
    sys_prompt = """
    Sei un Senior Event Manager esperto in Team Building.
    Analizza la RICHIESTA dell'utente e trova nel CATALOGO i format più pertinenti.
    
    THINKING PROCESS:
    1. Astrai la richiesta (es. "ponte tibetano" -> "Outdoor/Avventura").
    2. Cerca per ASSOCIAZIONE DI IDEE.
    3. Restituisci i NOMI DEI FORMAT (ID colonna 1).
    
    Output: SOLO lista Python. Es: ['Format A', 'Format B'].
    """
    
    model = genai.GenerativeModel(
        model_name=SEARCH_MODEL,
        generation_config={"temperature": 0.1},
        system_instruction=sys_prompt
    )
    try:
        response = model.generate_content(f"CATALOGO:\n{context_str}\n\nRICHIESTA UTENTE: {query}")
        if response.usage_metadata:
            st.session_state['token_usage']['total'] += response.usage_metadata.total_token_count
        
        match = re.search(r"(\[.*\])", response.text.strip(), re.DOTALL)
        return ast.literal_eval(match.group(1)) if match else []
    except: return []


# ==========================================
#              SIDEBAR CONTROL
# ==========================================
with st.sidebar:
    st.title("🦁 Manager")
    st.markdown("---")
    
    # 1. UPLOAD
    st.subheader("1. 📂 Carica File")
    uploaded_file = st.file_uploader("Trascina PDF o PPTX", type=['pdf', 'pptx', 'ppt'], label_visibility="collapsed")

    # LOGICA AUTO-ANALISI
    if uploaded_file:
        file_id = f"{uploaded_file.name}_{uploaded_file.size}"
        if st.session_state['last_processed_file'] != file_id:
            with st.spinner("⚡ Analisi Gemini 3.0..."):
                raw_text = read_file_content(uploaded_file)
                st.session_state['debug_raw_text'] = raw_text 
                
                if len(raw_text) > 10:
                    extracted = analyze_document_with_gemini(raw_text, [id_col] + cols)
                    if isinstance(extracted, list): extracted = extracted[0] if extracted else {}
                    if not isinstance(extracted, dict): extracted = {}

                    # FUZZY MATCH
                    extracted_name = str(extracted.get(id_col, "")).strip()
                    matches = difflib.get_close_matches(extracted_name, product_ids, n=1, cutoff=0.85)
                    
                    if matches:
                        existing_id = matches[0]
                        st.toast(f"Trovato: {existing_id}", icon="🔄")
                        
                        current_data = df.loc[existing_id].to_dict()
                        
                        # Trova nomi colonne
                        desc_col_name = next((c for c in cols if "descrizione" in c.lower()), "Descrizione Breve")
                        log_col_name = next((c for c in cols if "logistica" in c.lower()), "Logistica")
                        
                        st.session_state['pending_duplicate'] = {
                            'id': existing_id,
                            'desc': {
                                'col': desc_col_name,
                                'new': extracted.get(desc_col_name, ""),
                                'old': current_data.get(desc_col_name, "")
                            },
                            'log': {
                                'col': log_col_name,
                                'new': extracted.get(log_col_name, ""),
                                'old': current_data.get(log_col_name, "")
                            }
                        }
                        # Forza la selezione del format trovato
                        st.session_state['force_selection'] = existing_id
                        st.session_state['draft_data'] = {}
                    else:
                        st.session_state['pending_duplicate'] = None
                        st.session_state['draft_data'] = extracted if extracted else {}
                        if st.session_state['draft_data']:
                            st.toast("Nuovo format estratto!", icon="✨")
                else:
                    st.error("File illeggibile.")
                st.session_state['last_processed_file'] = file_id
    
    st.markdown("---")

    # 2. RICERCA
    st.subheader("2. 🔎 Cerca (AI)")
    q = st.text_input("Es. cucina, outdoor...", label_visibility="collapsed")
    if st.button("Cerca Format", use_container_width=True):
        if q:
            with st.spinner("Ricerca..."):
                res = search_ai(q, df)
                st.session_state['search_results'] = [x for x in res if x in product_ids] if res else []

    # RISULTATI RICERCA
    if st.session_state['search_results']:
        st.success(f"Trovati: {len(st.session_state['search_results'])}")
        
        desc_key = "Descrizione Breve"
        for c in cols:
            if "descrizione" in c.lower(): desc_key = c; break
            
        for rid in st.session_state['search_results']:
            # Card style in sidebar
            row_data = df.loc[rid]
            preview = str(row_data.get(desc_key, ""))[:60] + "..."
            
            with st.container():
                st.markdown(f"""
                <div class="result-card">
                    <div class="result-title">{rid}</div>
                    <div class="result-preview">{preview}</div>
                </div>
                """, unsafe_allow_html=True)
                if st.button("✏️ Modifica", key=f"btn_side_{rid}", use_container_width=True):
                    st.session_state['force_selection'] = rid
                    st.rerun()
        
        if st.button("❌ Reset Ricerca", use_container_width=True):
            st.session_state['search_results'] = None
            st.rerun()

    st.markdown("---")

    # 3. SELEZIONE MANUALE
    st.subheader("3. 📝 Selezione")
    all_options = [""] + sorted(product_ids)
    
    idx_sel = 0
    if st.session_state['force_selection']:
        if st.session_state['force_selection'] in all_options:
            idx_sel = all_options.index(st.session_state['force_selection'])
        st.session_state['force_selection'] = None
    
    selected_id = st.selectbox("Scegli Format:", all_options, index=idx_sel, label_visibility="collapsed")
    
    is_new_mode = False
    if st.session_state['draft_data'] and not st.session_state['pending_duplicate']:
        is_new_mode = True
        st.info("✏️ **CREAZIONE NUOVO**")
        if st.button("Annulla Creazione", use_container_width=True):
            st.session_state['draft_data'] = {}
            st.session_state['pending_changes'] = None
            st.session_state['last_processed_file'] = None
            st.rerun()


# ==========================================
#              MAIN COLUMN
# ==========================================

st.title("🦁 MasterTb Manager")

# 1. BOX GESTIONE DUPLICATI
if st.session_state['pending_duplicate']:
    dup_data = st.session_state['pending_duplicate']
    dup_id = dup_data['id']
    
    # Recupera dati desc
    d_col = dup_data['desc']['col']
    d_new = dup_data['desc']['new']
    d_old = dup_data['desc']['old']

    # Recupera dati log
    l_col = dup_data['log']['col']
    l_new = dup_data['log']['new']
    l_old = dup_data['log']['old']

    with st.container():
        st.warning(f"⚠️ **RILEVATO FORMAT ESISTENTE: '{dup_id}'**")
        st.markdown(f"L'AI propone di aggiornare **Descrizione** e **Logistica**. Modifica se necessario.")
        
        # COLONNA 1: DESCRIZIONE
        st.subheader(f"1. {d_col}")
        col_d1, col_d2 = st.columns(2)
        with col_d1:
            st.caption("🔴 Attuale")
            st.info(d_old if d_old else "(Vuoto)", icon="ℹ️")
        with col_d2:
            st.caption("🟢 Nuova (Editabile)")
            edited_desc = st.text_area("Modifica Nuova Descrizione", value=d_new, height=150, key="edit_d", label_visibility="collapsed")
        
        st.markdown("---")

        # COLONNA 2: LOGISTICA
        st.subheader(f"2. {l_col}")
        col_l1, col_l2 = st.columns(2)
        with col_l1:
            st.caption("🔴 Attuale")
            st.info(l_old if l_old else "(Vuoto)", icon="ℹ️")
        with col_l2:
            st.caption("🟢 Nuova (Editabile)")
            edited_log = st.text_area("Modifica Nuova Logistica", value=l_new, height=150, key="edit_l", label_visibility="collapsed")

        st.markdown("---")

        # AZIONI
        b1, b2 = st.columns([1, 4])
        with b1:
            if st.button("🔄 AGGIORNA ENTRAMBI", type="primary", use_container_width=True):
                try:
                    r_idx = product_ids.index(dup_id) + 2 
                    
                    # Update Descrizione
                    c_idx_d = cols.index(d_col) + 2
                    ws.update_cell(r_idx, c_idx_d, edited_desc)
                    
                    # Update Logistica
                    c_idx_l = cols.index(l_col) + 2
                    ws.update_cell(r_idx, c_idx_l, edited_log)

                    st.toast("Aggiornato!", icon="✅")
                    st.session_state['pending_duplicate'] = None
                    load_data.clear()
                    st.rerun()
                except Exception as e: st.error(f"Errore: {e}")
        with b2:
            if st.button("❌ IGNORA", use_container_width=True):
                st.session_state['pending_duplicate'] = None
                st.rerun()
    st.divider()


# 2. LOGICA VISIBILITÀ FORM
show_form = False
if is_new_mode:
    show_form = True
elif selected_id:
    show_form = True

if not show_form:
    st.info("👈 Usa la barra laterale per caricare un file, cercare o selezionare un format.")
    st.stop()


# 3. FORM PRINCIPALE (SCHEDA COMPLETA)
st.markdown("### 📝 Dettagli Format")

if is_new_mode:
    source_data = st.session_state['draft_data']
    current_id_val = str(source_data.get(id_col, ""))
    submit_label = "🧐 VERIFICA DATI (Step 1/2)"
else:
    source_data = df.loc[selected_id].to_dict()
    current_id_val = selected_id
    submit_label = "🧐 VERIFICA MODIFICHE (Step 1/2)"

with st.form("master_form"):
    form_values = {}
    
    # ID (Unico)
    if is_new_mode:
        new_id = st.text_input(f"**{id_col} (UNICO)**", value=current_id_val)
    else:
        st.text_input(f"**{id_col}**", value=current_id_val, disabled=True)
        new_id = current_id_val

    # Render dinamico campi
    for c in cols:
        val = str(source_data.get(c, ""))
        c_lower = c.lower()
        if "[[RIEMPIMENTO MANUALE]]" in val: val = ""
        
        # --- REGOLA NOVITÀ PER NUOVI FORMAT ---
        # Se è nuovo, la novità è SI per definizione.
        if is_new_mode and ("novità" in c_lower or "novita" in c_lower):
            val = "SI"
        
        # --- LOGICA WIDGET E LABELS ---
        
        # 1. METODO DI CALCOLO
        if "metodo" in c_lower and "calcolo" in c_lower:
            options_metodo = ["Standard", "Flat"]
            idx_metodo = 1 if "flat" in val.lower() else 0
            form_values[c] = st.selectbox(f"**{c}**", options_metodo, index=idx_metodo)

        # 2. SOCIAL / NOVITÀ (SI/NO)
        elif "social" in c_lower or "novità" in c_lower or "novita" in c_lower:
            options_bool = ["NO", "SI"]
            idx_bool = 1 if ("si" in val.lower() or "yes" in val.lower()) else 0
            form_values[c] = st.selectbox(f"**{c}**", options_bool, index=idx_bool)
        
        # 3. RANKING (1-5)
        elif "ranking" in c_lower:
            options_ranking = ["1", "2", "3", "4", "5"]
            try:
                curr_rank = str(int(float(val))) if val.strip() else "3"
                if curr_rank not in options_ranking: curr_rank = "3"
            except: curr_rank = "3"
            form_values[c] = st.selectbox(f"**{c}**", options_ranking, index=options_ranking.index(curr_rank))
            
        # 4. LINK AUTOMATICI E OBBLIGATORI
        elif "link" in c_lower:
            label = f"**{c}**"
            is_pdf_ppt = "pdf" in c_lower or "ppt" in c_lower
            
            # Se è un link PDF/PPT, aggiungi etichetta OBBLIGATORIO in rosso
            if is_pdf_ppt:
                label += " :red[(OBBLIGATORIO)]"
            
            slug = create_slug(new_id)
            
            # Website
            if "website" in c_lower:
                if is_new_mode or not val: val = f"https://www.teambuilding.it/project/{slug}/"
            
            # FILE PDF/PPT
            elif is_pdf_ppt:
                base_url = "https://teambuilding.it/preventivi/schede"
                lang = "eng" if "eng" in c_lower else "ita"
                ext = "pptx" if "ppt" in c_lower else "pdf"
                
                if is_new_mode or not val:
                    val = f"{base_url}/{lang}/{slug}.{ext}"
            
            form_values[c] = st.text_input(label, value=val)
            
        # 5. DURATA e ALTRI
        elif "durata" in c_lower and "ideale" in c_lower:
             form_values[c] = st.text_input(f"**{c}** (Media in ore)", value=val)

        else:
            height = 150 if "descrizione" in c_lower else 0
            if len(val) > 50 or height > 0:
                form_values[c] = st.text_area(f"**{c}**", value=val, height=height if height else None)
            else:
                form_values[c] = st.text_input(f"**{c}**", value=val)

    submitted = st.form_submit_button(submit_label, type="primary")

    if submitted:
        # --- VALIDAZIONE BLOCCANTE LINK TASSATIVI ---
        errors = []
        for c, val in form_values.items():
            c_lower = c.lower()
            if "link" in c_lower and ("pdf" in c_lower or "ppt" in c_lower):
                if not val.strip():
                    errors.append(f"Il campo '{c}' è TASSATIVO e non può essere vuoto!")
        
        if errors:
            for e in errors: st.error(e)
        else:
            # Se la validazione passa, procedi al calcolo modifiche
            changes = {}
            if is_new_mode:
                changes = {'_NEW_': True, 'id': new_id, 'data': form_values}
            else:
                for k, v in form_values.items():
                    original = str(source_data.get(k, ""))
                    if v != original:
                        changes[k] = {'old': original, 'new': v}
            st.session_state['pending_changes'] = changes

# 4. CONFERMA (DIFF VIEW)
if st.session_state['pending_changes']:
    st.divider()
    changes = st.session_state['pending_changes']
    
    if is_new_mode:
        st.info(f"✨ **STO CREANDO IL NUOVO FORMAT:** {changes['id']}")
        if st.button("✅ CONFERMA E SCRIVI (Definitivo)", type="primary"):
            if not changes['id'].strip():
                st.error("Manca ID!")
            elif changes['id'] in product_ids:
                st.error("Nome già esistente!")
            else:
                try:
                    row_to_append = [changes['id']] + [changes['data'][c] for c in cols]
                    ws.append_row(row_to_append)
                    st.success("Salvato!")
                    st.session_state['draft_data'] = {}
                    st.session_state['pending_changes'] = None
                    st.session_state['last_processed_file'] = None
                    load_data.clear()
                    st.rerun()
                except Exception as e: st.error(f"Errore: {e}")

    else:
        # EDIT MODE
        if not changes:
            st.success("✅ Nessuna modifica rilevata.")
            if st.button("Chiudi"):
                st.session_state['pending_changes'] = None
                st.rerun()
        else:
            st.warning("⚠️ **Rilevate Modifiche!** Controlla prima di salvare.")
            
            for k, v in changes.items():
                st.markdown(f"""
                <div class="diff-box">
                    <strong>{k}</strong><br>
                    <span class="old-val">OLD: {v['old'] if v['old'] else '(vuoto)'}</span> 
                    &nbsp;➡️&nbsp; 
                    <span class="new-val">NEW: {v['new']}</span>
                </div>
                """, unsafe_allow_html=True)
            
            c_yes, c_no = st.columns(2)
            with c_yes:
                if st.button("✅ CONFERMA SALVATAGGIO", type="primary"):
                    try:
                        row_idx = product_ids.index(selected_id) + 2
                        updates_count = 0
                        for col_name, val_dict in changes.items():
                            col_idx = cols.index(col_name) + 2
                            ws.update_cell(row_idx, col_idx, val_dict['new'])
                            updates_count += 1
                        st.success(f"Salvato! {updates_count} campi aggiornati.")
                        st.session_state['pending_changes'] = None
                        load_data.clear()
                        st.rerun()
                    except Exception as e: st.error(f"Errore: {e}")
            with c_no:
                if st.button("❌ Annulla"):
                    st.session_state['pending_changes'] = None
                    st.rerun()
